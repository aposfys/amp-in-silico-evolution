"""
PFES-AMP presentation — correct naming and logic
Run A: MACREL + ESM3
Run B: MACREL + ESM3 + PFES
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = "amps_pfes_macrel.pptx"

# ── Palette ────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x37, 0x5A)
DBLUE  = RGBColor(0x1A, 0x5C, 0x9A)   # Run A colour
TEAL   = RGBColor(0x00, 0x79, 0x6B)   # Run B colour
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DGREY  = RGBColor(0x22, 0x22, 0x22)
MGREY  = RGBColor(0x66, 0x66, 0x66)
LGREY  = RGBColor(0xF2, 0xF2, 0xF2)
BGREY  = RGBColor(0xE0, 0xE0, 0xE0)
LBLUE  = RGBColor(0xDB, 0xEA, 0xF8)
LGREEN = RGBColor(0xE0, 0xF2, 0xE8)
LYELL  = RGBColor(0xFD, 0xF6, 0xDD)
LORANG = RGBColor(0xFD, 0xED, 0xD8)
LPURP  = RGBColor(0xF3, 0xE5, 0xF5)
RED    = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xD4, 0x60, 0x10)

LABEL_A = "MACREL + ESM3"
LABEL_B = "MACREL + ESM3 + PFES"

W, H = 13.33, 7.5

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]


# ── Primitives ──────────────────────────────────────────────────────────────
def R(slide, x, y, w, h, fill=None, border=None, bw=0.75):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if border: s.line.color.rgb = border; s.line.width = Pt(bw)
    else: s.line.fill.background()
    return s


def T(slide, text, x, y, w, h, size=12, bold=False, italic=False,
      color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color if color else DGREY
    return tb


def I(slide, path, x, y, w, h=None):
    if not os.path.exists(path): return
    if h: slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else: slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


def footer_bar(slide):
    R(slide, 0, H-0.48, W, 0.48, fill=LGREY)
    R(slide, 0, H-0.48, W, 0.04, fill=BGREY)
    T(slide,
      "National & Kapodistrian University of Athens   ·   Hellenic Pasteur Institute   ·   Evolutionary Genomics Group",
      0.3, H-0.43, W-0.6, 0.35, size=8.5, color=MGREY, align=PP_ALIGN.CENTER)


def slide_title(slide, title, sub=""):
    T(slide, title, 0.35, 0.14, W-0.7, 0.55, size=26, bold=True, color=NAVY)
    if sub:
        T(slide, sub, 0.35, 0.68, W-0.7, 0.28, size=11, italic=True, color=MGREY)
    R(slide, 0.35, 0.96 if sub else 0.72, W-0.7, 0.04, fill=BGREY)


def col_labels(slide, left, right, lc=DBLUE, rc=TEAL):
    R(slide, 0.3, 0.98, 6.15, 0.36, fill=lc)
    R(slide, 6.68, 0.98, 6.3, 0.36, fill=rc)
    T(slide, left,  0.45, 1.02, 5.8, 0.28, size=12, bold=True, color=WHITE)
    T(slide, right, 6.83, 1.02, 6.0, 0.28, size=12, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
T(sl, "Investigating novel AMPs using\nmachine learning, structure prediction\nand in silico evolution",
  0.9, 1.2, W-1.8, 2.9, size=38, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
T(sl, "MSc Bioinformatics – Computational Biology\nNational & Kapodistrian University of Athens\nApostolos Fysekidis",
  0.9, 4.3, W-1.8, 1.4, size=16, color=MGREY, align=PP_ALIGN.CENTER)
R(sl, 3.4, 6.0, 2.9, 0.07, fill=DBLUE)
R(sl, 7.0, 6.0, 2.9, 0.07, fill=TEAL)
footer_bar(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The PFES loop (original framework)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
slide_title(sl, "The PFES Framework", "Original method — Sahakyanhk et al., PNAS 2022")

bx, bw, bh, gap = 0.28, 2.96, 4.0, 0.14
by = 1.08
colors  = [LBLUE, LGREEN, LYELL, LORANG]
borders = [DBLUE, TEAL, RGBColor(0xD4,0xA0,0x17), ORANGE]
heads   = ["1. MUTATE", "2. PREDICT", "3. SCORE", "4. SELECT"]
bodies  = [
    "Each of the 100 sequences in the population is mutated once per generation.\n\nOperators: substitution (most frequent), insertion (+), deletion (−), scramble (*), duplication (d, rare).",
    "ESMFold v1 predicts the 3D structure of every mutated sequence.\n\nOutputs per sequence:\n• pLDDT — per-residue confidence [0–1]\n• pTM — global fold quality [0–1]\n• Backbone coordinates → contacts, secondary structure",
    "A multiplicative fitness score is computed from structural quality terms only.\n\nAll terms in [0,1]. Any near-zero term collapses the full score — hard selection pressure on every dimension simultaneously.",
    "100 survivors are sampled from the pool of 200 (old + new) using Boltzmann-weighted probabilities.\n\nβ = 20: sequences with higher scores are exponentially more likely to survive, while some diversity is retained.",
]
for i in range(4):
    cx = bx + i*(bw+gap)
    R(sl, cx, by, bw, bh, fill=colors[i], border=borders[i])
    T(sl, heads[i], cx+0.14, by+0.12, bw-0.28, 0.38, size=14, bold=True, color=borders[i])
    R(sl, cx+0.14, by+0.54, bw-0.28, 0.03, fill=BGREY)
    T(sl, bodies[i], cx+0.14, by+0.64, bw-0.28, bh-0.78, size=10.5, color=DGREY, wrap=True)
    if i < 3:
        T(sl, "→", cx+bw+0.01, by+bh/2-0.2, 0.16, 0.4, size=18, color=MGREY, align=PP_ALIGN.CENTER)

R(sl, 0.28, by+bh+0.16, W-0.56, 0.42, fill=RGBColor(0xE8,0xEA,0xF6), border=BGREY)
T(sl, "500 generations × 100 sequences = 50,000 structure predictions per run   ·   Your project replaced ESMFold v1 with ESM3",
  0.42, by+bh+0.22, W-0.84, 0.3, size=10, color=NAVY, align=PP_ALIGN.CENTER)
footer_bar(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Original PFES score formula
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
slide_title(sl, "Original PFES Fitness Score", "Purely structural — no AMP scoring")

# Formula
R(sl, 0.28, 0.88, W-0.56, 0.58, fill=RGBColor(0xE8,0xEA,0xF6), border=BGREY)
T(sl, "S  =  pLDDT  ×  pTM  ×  contact_density  ×  len_penalty  ×  helix_penalty  ×  β_penalty",
  0.42, 0.96, W-0.84, 0.42, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Note: single chain simplification
R(sl, 0.28, 1.5, W-0.56, 0.34, fill=LGREY)
T(sl, "For single-chain peptides: ipLDDT = 1.0 and inter-contact term = 1.0 — both cancel out, leaving the 6 terms above.",
  0.42, 1.55, W-0.84, 0.24, size=9.5, color=MGREY, italic=True, align=PP_ALIGN.CENTER)

# 6 term boxes
terms = [
    ("pLDDT",         "from ESM3",       "Per-residue confidence.\n>0.9 = well-structured\n<0.5 = likely disordered\nMean over all residues.",              DBLUE, LBLUE),
    ("pTM",           "from ESM3",       "Global fold quality.\n>0.5 = reasonable topology\n>0.8 = high confidence\nSensitive to overall fold.",             TEAL,  LGREEN),
    ("contact\ndensity","from ESM3 PDB", "(Cβ–Cβ contacts within 6Å\n+ seq_len) / seq_len\nRewards compact globular\nfolds.",                               RGBColor(0xD4,0xA0,0x17), LYELL),
    ("len_penalty",   "PFES original",   "1 − σ(len, 250, 0.2)\nThreshold 250 aa.\nPenalises sequences that\ngrow beyond typical size.",                    ORANGE, LORANG),
    ("helix_penalty", "PFES original",   "1 − σ(max_helix, 20, 0.5)\nThreshold 20 aa.\nPenalises any single helix\nlonger than 20 residues.",               RGBColor(0x7B,0x1F,0xA2), LPURP),
    ("β_penalty",     "PFES original",   "1 − σ(max_β, 12, 0.6)\nThreshold 12 aa.\nPenalises any β-strand\nlonger than 12 residues.",                       RGBColor(0x4A,0x14,0x8C), RGBColor(0xEA,0xE0,0xF8)),
]

tw = (W - 0.7) / 6
for i, (name, src, body, tc, bg) in enumerate(terms):
    cx = 0.28 + i * (tw + 0.06)
    R(sl, cx, 1.9, tw, 4.5, fill=bg, border=BGREY)
    T(sl, name, cx+0.1, 1.98, tw-0.2, 0.52, size=13, bold=True, color=tc, align=PP_ALIGN.CENTER)
    R(sl, cx+0.1, 2.54, tw-0.2, 0.26, fill=tc)
    T(sl, src,  cx+0.1, 2.56, tw-0.2, 0.22, size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, body, cx+0.1, 2.88, tw-0.2, 3.4,  size=10, color=DGREY, wrap=True)
    if i < 5:
        T(sl, "×", cx+tw+0.01, 3.8, 0.08, 0.3, size=14, color=MGREY, align=PP_ALIGN.CENTER)

footer_bar(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MACREL: what it adds
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
slide_title(sl, "MACREL — What It Adds to the Score")

# Left: what MACREL is
R(sl, 0.28, 0.88, 6.2, 5.6, fill=LBLUE, border=BGREY)
T(sl, "What MACREL provides", 0.46, 0.96, 5.8, 0.36, size=14, bold=True, color=DBLUE)

rows = [
    ("Input",    "Any amino acid sequence — no 3D structure required"),
    ("Model",    "Two separate ONNX classifiers trained on physicochemical features:\n23 features including charge, hydrophobicity, hydrophobic moment,\namino acid group composition, Boman index, instability index"),
    ("Output 1", "AMP probability [0–1]\nHow likely the sequence is a true antimicrobial peptide\n> 0.5 → classified as AMP"),
    ("Output 2", "Hemolytic probability [0–1]\nHow likely the sequence lyses red blood cells\nTrainned on database peptides — unreliable for highly cationic\nsequences evolved by PFES (outputs ~0 for all cationic AMPs)"),
    ("Usage",    "Called as a subprocess each generation:\nmacrel peptides --fasta batch.faa\nTakes a FASTA of the 100 mutated sequences → returns predictions"),
]
for i, (label, body) in enumerate(rows):
    y = 1.44 + i * 0.84
    R(sl, 0.42, y, 1.35, 0.68, fill=NAVY)
    T(sl, label, 0.46, y+0.14, 1.27, 0.42, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, body, 1.87, y+0.06, 4.45, 0.74, size=10, color=DGREY, wrap=True)

# Right: hemolytic problem and solution
R(sl, 6.68, 0.88, 6.3, 2.6, fill=LORANG, border=BGREY)
T(sl, "⚠  Hemolytic probability problem", 6.86, 0.96, 6.0, 0.36, size=13, bold=True, color=ORANGE)
T(sl,
  "MACREL's hemolytic classifier was trained on known database peptides. "
  "The sequences PFES evolves are highly cationic amphipathic helices — outside the training distribution. "
  "MACREL outputs hemolytic probability ≈ 0.001 for all of them, providing zero selection pressure against toxic peptides.",
  6.86, 1.4, 6.0, 1.95, size=10.5, color=DGREY, wrap=True)

R(sl, 6.68, 3.58, 6.3, 2.9, fill=LGREEN, border=BGREY)
T(sl, "✓  Biophysical proxy (Option 2)", 6.86, 3.66, 6.0, 0.36, size=13, bold=True, color=TEAL)
T(sl, "No dataset needed — based on known biology:", 6.86, 4.1, 6.0, 0.28, size=11, bold=True, color=NAVY)

R(sl, 6.78, 4.44, 6.1, 0.56, fill=WHITE, border=BGREY)
T(sl, "hemo_proxy = sigmoid( hydro_ratio × 10 − min(charge, 8) × 0.5 − 2 )",
  6.86, 4.5, 6.0, 0.44, size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

T(sl,
  "Hemolysis requires hydrophobicity (membrane insertion). "
  "High positive charge reduces hemolytic risk: cationic peptides preferentially bind anionic bacterial membranes over zwitterionic mammalian ones. "
  "Charge is capped at +8 to prevent extreme K/R accumulation from eliminating the penalty entirely.",
  6.86, 5.08, 6.0, 1.25, size=10, color=DGREY, wrap=True)

footer_bar(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Two approaches: score formulas
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
slide_title(sl, "Two Experimental Runs")
col_labels(sl, LABEL_A, LABEL_B)

# Left formula
R(sl, 0.28, 1.4, 6.2, 0.96, fill=LBLUE, border=DBLUE, bw=1.5)
T(sl, "Score  =  pLDDT  ×  pTM\n            ×  AMP_prob  ×  (1 − hemo_proxy)",
  0.44, 1.48, 5.9, 0.82, size=13, bold=True, color=DBLUE, align=PP_ALIGN.CENTER)

# Right formula
R(sl, 6.68, 1.4, 6.3, 1.14, fill=LGREEN, border=TEAL, bw=1.5)
T(sl, "Score  =  pLDDT  ×  pTM  ×  contact_density\n            ×  len_penalty  ×  helix_penalty  ×  β_penalty\n            ×  AMP_prob  ×  (1 − hemo_proxy)",
  6.84, 1.48, 6.0, 1.0, size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Where each term comes from
def term_origin(sl, x, y, w, items):
    for i, (term, src, color) in enumerate(items):
        R(sl, x, y+i*0.46, w, 0.42, fill=LGREY if i%2==0 else WHITE)
        T(sl, term, x+0.1, y+i*0.46+0.07, 2.0, 0.3, size=11, bold=True, color=color)
        T(sl, src,  x+2.2, y+i*0.46+0.07, w-2.3, 0.3, size=10.5, color=DGREY)

left_items = [
    ("pLDDT",        "ESM3 structure prediction — per-residue confidence",         DBLUE),
    ("pTM",          "ESM3 structure prediction — global fold quality",             DBLUE),
    ("AMP_prob",     "MACREL ML classifier — AMP likelihood from sequence",        RGBColor(0x2E,0x7D,0x32)),
    ("1 − hemo",     "Biophysical proxy — hemolytic safety from sequence",         ORANGE),
]
right_items = [
    ("pLDDT",        "ESM3 — same as left",                                        DBLUE),
    ("pTM",          "ESM3 — same as left",                                        DBLUE),
    ("contact_density","PFES original — compactness from ESM3 PDB output",        RGBColor(0xD4,0xA0,0x17)),
    ("len_penalty",  "PFES original — penalises sequences > 30 aa",               ORANGE),
    ("helix_penalty","PFES original — penalises single helix > 20 aa",            RGBColor(0x7B,0x1F,0xA2)),
    ("β_penalty",    "PFES original — penalises β-strand > 12 aa",                RGBColor(0x4A,0x14,0x8C)),
    ("AMP_prob",     "MACREL — same as left",                                      RGBColor(0x2E,0x7D,0x32)),
    ("1 − hemo",     "Biophysical proxy — same as left",                           ORANGE),
]

T(sl, "Where each term comes from:", 0.3, 2.46, 5.9, 0.28, size=10, bold=True, color=MGREY)
term_origin(sl, 0.28, 2.76, 6.2, left_items)

T(sl, "Where each term comes from:", 6.68, 2.64, 6.1, 0.28, size=10, bold=True, color=MGREY)
term_origin(sl, 6.68, 2.94, 6.3, right_items)

# Bottom: what the PFES terms add
R(sl, 0.28, H-1.3, W-0.56, 0.74, fill=RGBColor(0xE8,0xEA,0xF6), border=BGREY)
T(sl, "What the PFES structural terms add:", 0.44, H-1.22, 3.2, 0.28, size=11, bold=True, color=NAVY)
T(sl, "Contact density rewards compact folds — prevents disordered blobs. "
      "Length penalty prevents the duplication operator from growing sequences to 70+ aa. "
      "Helix and β penalties prevent secondary structure run-away. "
      "Without them, the optimiser exploits duplication: a 24 aa sequence scoring 0.6 "
      "becomes a 48 aa duplicate scoring ~0.8 — an artefact, not a better peptide.",
  3.7, H-1.22, W-4.0, 0.62, size=10, color=DGREY, wrap=True)

footer_bar(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Run parameters
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
slide_title(sl, "Run Parameters")
col_labels(sl, LABEL_A, LABEL_B)

params = [
    ("Structure predictor",   "ESM3",                  "ESM3",                  "Replaced ESMFold v1 from original PFES"),
    ("Selection mode",        "Weak  (β = 20)",         "Weak  (β = 20)",         "Boltzmann-weighted sampling"),
    ("Population size",       "100",                   "100",                   "Sequences per generation"),
    ("Generations",           "500",                   "500",                   ""),
    ("Starting length",       "24 aa  (random)",       "24 aa  (random)",       ""),
    ("Repeat filter",         "ON  (--norepeat)",       "ON  (--norepeat)",       "No sequence evaluated twice"),
    ("AMP scoring",           "MACREL ML model",       "MACREL ML model",       "Probability > 0.5 → AMP"),
    ("Hemolytic scoring",     "Biophysical proxy",     "Biophysical proxy",     "sigmoid(hydro×10 − min(charge,8)×0.5 − 2)"),
    ("Contact density",       "✗  not in score",       "✓  in score",           "Key PFES term"),
    ("Length penalty",        "✗  not in score",       "✓  threshold ~30 aa",   "Key PFES term"),
    ("Helix penalty",         "✗  not in score",       "✓  threshold 20 aa",    "Key PFES term"),
    ("β-strand penalty",      "✗  not in score",       "✓  threshold 12 aa",    "Key PFES term"),
]

R(sl, 0.28, 1.4, W-0.56, 0.42, fill=NAVY)
for label, cx, cw in [("Parameter",0.38,4.2),("MACREL + ESM3",4.68,3.4),("MACREL + ESM3 + PFES",8.18,3.4),("Notes",11.68,1.5)]:
    T(sl, label, cx, 1.46, cw, 0.3, size=10.5, bold=True, color=WHITE)

for i, (param, va, vb, note) in enumerate(params):
    y = 1.86 + i*0.38
    bg = LGREY if i%2==0 else WHITE
    R(sl, 0.28, y, W-0.56, 0.36, fill=bg)
    T(sl, param, 0.38, y+0.05, 4.2, 0.28, size=10, color=DGREY)
    diff = va != vb
    ca = DBLUE if (diff and "✗" not in va and "✓" not in va) else (RED if "✗" in va else (TEAL if "✓" in va else DGREY))
    cb = TEAL  if (diff and "✓" in vb) else (RED if "✗" in vb else DGREY)
    T(sl, va, 4.68, y+0.05, 3.4, 0.28, size=10, bold=diff, color=ca)
    T(sl, vb, 8.18, y+0.05, 3.4, 0.28, size=10, bold=diff, color=cb)
    if note:
        T(sl, note, 11.68, y+0.05, 1.5, 0.28, size=8, color=MGREY, wrap=True)

footer_bar(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MACREL + ESM3 results
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
slide_title(sl, f"{LABEL_A} — Results")
R(sl, 0.28, 0.88, 3.7, 0.36, fill=DBLUE)
T(sl, LABEL_A, 0.42, 0.92, 3.5, 0.28, size=11, bold=True, color=WHITE)

# Best candidate
R(sl, 0.28, 1.32, 4.42, 3.22, fill=LBLUE, border=BGREY)
T(sl, "Best candidate  (gen 482)", 0.44, 1.4, 4.1, 0.32, size=12, bold=True, color=NAVY)
T(sl, "VWMERKMIAQKKIRKLQMIKMKLQMIKMKRHIRQAK\nRMIRQAKIRQLQVTIAMAKKMRKKLEEMMDFLFKMWRWADIH",
  0.44, 1.8, 4.1, 0.65, size=9.5, bold=True, color=DBLUE, wrap=True)

stats_a = [("Score","0.830"),("AMP prob","0.999"),("Hemo proxy","0.001"),
           ("pLDDT","0.990"),("pTM","0.840"),("Length","78 aa  ⚠")]
for i,(k,v) in enumerate(stats_a):
    row,col = i//2, i%2
    y = 2.56 + row*0.38
    x = 0.42 + col*2.1
    T(sl, f"{k}:", x, y, 1.0, 0.3, size=10, color=MGREY)
    c = RED if "⚠" in v else DBLUE
    T(sl, v, x+1.05, y, 1.0, 0.3, size=10, bold=True, color=c)

R(sl, 0.28, 3.72, 4.42, 0.58, fill=LORANG, border=BGREY)
T(sl, "⚠  Sequences grew to 78 aa via duplication — no length penalty in score",
  0.44, 3.8, 4.2, 0.44, size=10, color=ORANGE, wrap=True)

R(sl, 0.28, 4.38, 4.42, 0.5, fill=LGREY, border=BGREY)
T(sl, "Final gen: mean score 0.784  ·  mean length ~71 aa  ·  mean hemo proxy 0.04",
  0.44, 4.46, 4.2, 0.36, size=10, color=DGREY, wrap=True)

# Plots
I(sl, "results_macrel_prod/analysis/Evolution.png",        4.9, 0.88, 4.1)
I(sl, "results_macrel_prod/analysis/Score_components.png", 9.1, 0.88, 4.1)
I(sl, "results_macrel_prod/analysis/AA_composition.png",   4.9, 3.5, 4.1)
I(sl, "results_macrel_prod/analysis/Score_distribution.png",9.1, 3.5, 4.1)
for lbl,x,y in [("Fitness over generations",4.9,0.86),("Score components",9.1,0.86),
                 ("Amino acid composition",4.9,3.48),("Score distribution",9.1,3.48)]:
    T(sl, lbl, x, y, 4.0, 0.2, size=8.5, italic=True, color=MGREY)
footer_bar(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MACREL + ESM3 + PFES results
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
slide_title(sl, f"{LABEL_B} — Results")
R(sl, 0.28, 0.88, 4.6, 0.36, fill=TEAL)
T(sl, LABEL_B, 0.42, 0.92, 4.4, 0.28, size=11, bold=True, color=WHITE)

# Best candidate
R(sl, 0.28, 1.32, 4.42, 3.22, fill=LGREEN, border=BGREY)
T(sl, "Best candidate  (gen 483)", 0.44, 1.4, 4.1, 0.32, size=12, bold=True, color=NAVY)
T(sl, "RRVYKKFFAPRVRLKRLAKAIKLVRK",
  0.44, 1.82, 4.1, 0.36, size=15, bold=True, color=TEAL)

stats_b = [("Score","0.724"),("AMP prob","0.993"),("Hemo proxy","0.047  ✓"),
           ("pLDDT","0.980"),("pTM","0.550"),("Length","26 aa  ✓")]
for i,(k,v) in enumerate(stats_b):
    row,col = i//2, i%2
    y = 2.28 + row*0.38
    x = 0.42 + col*2.1
    T(sl, f"{k}:", x, y, 1.0, 0.3, size=10, color=MGREY)
    c = TEAL if "✓" in v else DGREY
    T(sl, v, x+1.05, y, 1.0, 0.3, size=10, bold=True, color=c)

R(sl, 0.28, 3.46, 4.42, 0.58, fill=LGREEN, border=BGREY)
T(sl, "✓  Sequences constrained to 24–27 aa — PFES length penalty active",
  0.44, 3.54, 4.2, 0.44, size=10, color=TEAL, wrap=True)

R(sl, 0.28, 4.12, 4.42, 0.5, fill=LGREEN, border=BGREY)
T(sl, "Final gen: mean score 0.672  ·  mean length 26 aa  ·  mean hemo proxy 0.062",
  0.44, 4.2, 4.2, 0.36, size=10, color=TEAL, wrap=True)

R(sl, 0.28, 4.7, 4.42, 0.48, fill=LGREY, border=BGREY)
T(sl, "Converged to short amphipathic α-helices — consistent with known AMP membrane-disruption mechanism",
  0.44, 4.77, 4.2, 0.36, size=9.5, color=DGREY, wrap=True)

# Plots
I(sl, "results_macrel_pfes_pro/analysis/Evolution.png",        4.9, 0.88, 4.1)
I(sl, "results_macrel_pfes_pro/analysis/Score_components.png", 9.1, 0.88, 4.1)
I(sl, "results_macrel_pfes_pro/analysis/AA_composition.png",   4.9, 3.5, 4.1)
I(sl, "results_macrel_pfes_pro/analysis/Score_distribution.png",9.1, 3.5, 4.1)
for lbl,x,y in [("Fitness over generations",4.9,0.86),("Score components",9.1,0.86),
                 ("Amino acid composition",4.9,3.48),("Score distribution",9.1,3.48)]:
    T(sl, lbl, x, y, 4.0, 0.2, size=8.5, italic=True, color=MGREY)
footer_bar(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Comparative plots
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
slide_title(sl, "Direct Comparison", f"{LABEL_A}  vs  {LABEL_B}")

R(sl, 0.28, 0.9, 2.6, 0.32, fill=DBLUE)
T(sl, LABEL_A, 0.38, 0.94, 2.45, 0.24, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
R(sl, 3.0, 0.9, 3.0, 0.32, fill=TEAL)
T(sl, LABEL_B, 3.1, 0.94, 2.85, 0.24, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

plots = [
    ("comparison_plots/1_fitness_over_generations.png", "Fitness over generations"),
    ("comparison_plots/2_score_components.png",         "Score components"),
    ("comparison_plots/3_aa_composition.png",           "Amino acid composition"),
    ("comparison_plots/4_score_distribution.png",       "Score distribution — final generation"),
]
positions = [(0.28,1.28),(6.7,1.28),(0.28,4.08),(6.7,4.08)]
for (path,lbl),(px,py) in zip(plots, positions):
    I(sl, path, px, py+0.22, 6.2)
    T(sl, lbl, px, py, 6.2, 0.22, size=9, italic=True, color=MGREY)
footer_bar(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Length artefact
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
slide_title(sl, "The Length Problem", "Why PFES structural penalties matter for AMP design")

I(sl, "comparison_plots/5_sequence_length.png", 0.28, 1.0, 7.6)

R(sl, 8.1, 1.0, 4.95, 5.5, fill=LGREY, border=BGREY)
T(sl, "What the plot shows", 8.28, 1.1, 4.6, 0.34, size=13, bold=True, color=NAVY)

points = [
    (DBLUE,  f"{LABEL_A}: sequences grow to 78 aa by generation ~300. "
             "The duplication operator (d) copies a well-scoring motif, roughly doubling most score terms. "
             "Without a length penalty in the score this is not discouraged — the optimiser exploits it."),
    (TEAL,   f"{LABEL_B}: sequences stay at 24–27 aa throughout all 500 generations. "
             "The length penalty (sigmoid centred at ~30 aa) makes any sequence longer than a typical AMP pay a score cost."),
    (ORANGE, "Clinical relevance: most therapeutic AMPs are 10–40 aa. Longer sequences face pharmacokinetic problems — "
             "poor membrane penetration, rapid proteolysis, higher manufacturing cost."),
    (NAVY,   f"Score artefact: {LABEL_A} best score 0.83 is inflated by length. "
             f"The shorter {LABEL_B} peptide RRVYKKFFAPRVRLKRLAKAIKLVRK at 0.72 is the more honest and drug-relevant result."),
]
for i,(color,text) in enumerate(points):
    R(sl, 8.2, 1.56+i*1.2, 0.1, 1.0, fill=color)
    T(sl, text, 8.42, 1.56+i*1.2, 4.52, 1.12, size=10, color=DGREY, wrap=True)

footer_bar(sl)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Summary
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
slide_title(sl, "Summary")
col_labels(sl, LABEL_A, LABEL_B)

rows = [
    ("Score formula",          "pLDDT × pTM × AMP × (1−hemo)",          "pLDDT × pTM × CD × penalties × AMP × (1−hemo)"),
    ("Structure predictor",    "ESM3",                                    "ESM3"),
    ("AMP classifier",         "MACREL",                                  "MACREL"),
    ("Hemolytic term",         "Biophysical proxy",                       "Biophysical proxy"),
    ("PFES structural terms",  "✗  absent",                               "✓  contact density + 3 penalties"),
    ("Best score",             "0.830  (inflated)",                       "0.724  (honest)"),
    ("Best sequence length",   "78 aa  ⚠  artefact",                     "26 aa  ✓  realistic AMP"),
    ("Best AMP probability",   "0.999",                                   "0.993"),
    ("Best hemo proxy",        "0.001  (low charge cap prevents it)",     "0.047  ✓"),
    ("Final pop mean score",   "0.784",                                   "0.672"),
    ("Viable drug candidate?", "Questionable — unrealistic length  ⚠",   "Yes — short, structured, low hemo  ✓"),
]

R(sl, 0.28, 1.4, W-0.56, 0.4, fill=NAVY)
for lbl,cx,cw in [("Metric",0.38,4.0),(LABEL_A,4.48,4.1),(LABEL_B,8.68,4.3)]:
    T(sl, lbl, cx, 1.46, cw, 0.3, size=10, bold=True, color=WHITE)

for i,(metric,va,vb) in enumerate(rows):
    y = 1.84 + i*0.42
    bg = LGREY if i%2==0 else WHITE
    R(sl, 0.28, y, W-0.56, 0.4, fill=bg)
    T(sl, metric, 0.38, y+0.06, 4.0, 0.3, size=10, color=DGREY)
    ca = RED  if "✗" in va or "⚠" in va else (DBLUE if "✓" in va else DGREY)
    cb = TEAL if "✓" in vb else (RED if "✗" in vb or "⚠" in vb else DGREY)
    T(sl, va, 4.48, y+0.06, 4.1, 0.3, size=10, bold=True, color=ca)
    T(sl, vb, 8.68, y+0.06, 4.3, 0.3, size=10, bold=True, color=cb)

R(sl, 0.28, H-1.18, W-0.56, 0.62, fill=RGBColor(0xE8,0xEA,0xF6), border=BGREY)
T(sl, "Conclusion:", 0.44, H-1.1, 1.4, 0.26, size=11, bold=True, color=NAVY)
T(sl, "Adding the original PFES structural penalties to MACREL + ESM3 prevents the length artefact, "
      "constrains sequences to realistic AMP sizes, and produces a score that reflects true peptide quality. "
      "The lower score (0.72 vs 0.83) is more meaningful, not worse.",
  1.9, H-1.1, W-2.2, 0.5, size=10.5, color=DGREY, wrap=True)

footer_bar(sl)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved → {OUT}  ({os.path.getsize(OUT)//1024} KB)")
