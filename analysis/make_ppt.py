"""
PPT: individual run plots + comparative plots.
Focus: clean look and feel, one graph per slide, full size.
"""
import os, warnings
import pandas as pd
import numpy as np
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mtick
from PIL import Image as PILImage
warnings.filterwarnings('ignore')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── config ────────────────────────────────────────────────────────────────────
LOG_A  = "results_macrel_300/progress.log"
LOG_B  = "results_macrel_pfes_300/progress.log"
DIR_A  = "results_macrel_300/analysis"
DIR_B  = "results_macrel_pfes_300/analysis"
CDIR   = "comparison_plots_300"
OUT    = "comparison_macrel_vs_pfes_300.pptx"
os.makedirs(CDIR, exist_ok=True)

CA_s, CB_s = "#1565C0", "#2E7D32"
CA_l, CB_l = "#90CAF9", "#A5D6A7"
LABEL_A, LABEL_B = "MACREL + ESM3", "MACREL + ESM3 + PFES"

# ══════════════════════════════════════════════════════════════════════════════
# STYLE — clean, publication-quality
# ══════════════════════════════════════════════════════════════════════════════
STYLE = {
    'font.family': 'DejaVu Sans',
    'font.size': 14,
    'axes.titlesize': 17,
    'axes.titleweight': 'bold',
    'axes.titlepad': 14,
    'axes.labelsize': 15,
    'axes.labelcolor': '#333333',
    'axes.edgecolor': '#CCCCCC',
    'axes.linewidth': 1.0,
    'axes.spines.top': False,
    'axes.spines.right': False,
    'axes.grid': True,
    'grid.color': '#E5E5E5',
    'grid.linewidth': 0.8,
    'grid.linestyle': '-',
    'xtick.labelsize': 13,
    'ytick.labelsize': 13,
    'xtick.color': '#555555',
    'ytick.color': '#555555',
    'legend.fontsize': 13,
    'legend.framealpha': 0.95,
    'legend.edgecolor': '#CCCCCC',
    'legend.borderpad': 0.6,
    'figure.facecolor': 'white',
    'axes.facecolor': 'white',
}
plt.rcParams.update(STYLE)
DPI = 200
FS  = (12, 6)      # standard figure size
FS_W = (13, 5.5)   # wide (AA composition)

# ══════════════════════════════════════════════════════════════════════════════
# DATA
# ══════════════════════════════════════════════════════════════════════════════
def load(path):
    df = pd.read_csv(path, sep='\t', comment='#', low_memory=False)
    amp = 'amp_prob' if 'amp_prob' in df.columns else 's_amp'
    df['amp'] = pd.to_numeric(df[amp], errors='coerce')
    for c in ['score','hemo_prob','mean_plddt','ptm','seq_len']:
        df[c] = pd.to_numeric(df[c], errors='coerce')
    df['gen'] = df['gndx'].str.extract(r'(\d+)').astype(float)
    return df.dropna(subset=['score','gen'])

A, B = load(LOG_A), load(LOG_B)

def sm(s, w=7): return s.rolling(w, min_periods=1, center=True).mean()
def gstat(df, col):
    g = df.groupby('gen')[col]
    return g.max(), g.mean(), g.quantile(0.25), g.quantile(0.75)

# ══════════════════════════════════════════════════════════════════════════════
# COMPARISON PLOTS
# ══════════════════════════════════════════════════════════════════════════════
def comp_plot(col, ylabel, title, fname, ylim=None, hline=None, pct=False):
    fig, ax = plt.subplots(figsize=FS)
    for df, c, cl, label in [(A,CA_s,CA_l,LABEL_A),(B,CB_s,CB_l,LABEL_B)]:
        mx, mn, q25, q75 = gstat(df, col)
        ax.fill_between(mx.index, sm(q25), sm(q75), color=c, alpha=0.12, lw=0)
        ax.plot(mx.index, sm(mx), color=c, lw=3.0,
                label=f'{label}  —  best / gen', solid_capstyle='round')
        ax.plot(mx.index, sm(mn), color=c, lw=1.5, ls='--', alpha=0.65,
                label=f'{label}  —  mean / gen', solid_capstyle='round')
    if hline:
        ax.axhline(hline[0], color='#888888', lw=1.8, ls=':', alpha=0.9)
        ax.text(298, hline[0]+0.5, hline[1], ha='right', va='bottom',
                fontsize=11, color='#888888', style='italic')
    if ylim: ax.set_ylim(*ylim)
    if pct: ax.yaxis.set_major_formatter(mtick.PercentFormatter(xmax=1))
    ax.set_xlabel('Generation'); ax.set_ylabel(ylabel)
    ax.set_title(title)
    ax.legend(loc='lower right', ncol=1)
    fig.tight_layout(pad=1.5)
    fig.savefig(f'{CDIR}/{fname}', dpi=DPI, bbox_inches='tight'); plt.close()

comp_plot('score',     'Fitness score',         'Fitness Over Generations',
          'c_fitness.png',    (0, 0.85))
comp_plot('seq_len',   'Sequence length (aa)',   'Sequence Length Over Generations',
          'c_length.png',
          hline=(30, '30 aa  ·  therapeutic target'))
comp_plot('amp',       'AMP probability',        'AMP Probability Over Generations',
          'c_amp.png',        (0, 1.05), pct=False)
comp_plot('hemo_prob', 'Hemolytic proxy',         'Hemolytic Proxy Over Generations',
          'c_hemo.png',       (0, 0.5))
comp_plot('mean_plddt','pLDDT',                  'pLDDT Over Generations',
          'c_plddt.png',      (0.4, 1.05))

# ── score distribution ────────────────────────────────────────────────────────
fa_s = A[A.gen==A.gen.max()].score.dropna()
fb_s = B[B.gen==B.gen.max()].score.dropna()
fig, ax = plt.subplots(figsize=FS)
bins = np.linspace(0, 0.85, 28)
ax.hist(fa_s, bins=bins, color=CA_s, alpha=0.70, density=True,
        label=f'{LABEL_A}   μ = {fa_s.mean():.3f}', zorder=3)
ax.hist(fb_s, bins=bins, color=CB_s, alpha=0.70, density=True,
        label=f'{LABEL_B}   μ = {fb_s.mean():.3f}', zorder=3)
ax.axvline(fa_s.mean(), color=CA_s, lw=2.5, ls='--', zorder=4)
ax.axvline(fb_s.mean(), color=CB_s, lw=2.5, ls='--', zorder=4)
ax.set_xlabel('Fitness score'); ax.set_ylabel('Density')
ax.set_title('Score Distribution — Final Generation (gen 299)')
ax.legend(loc='upper left')
fig.tight_layout(pad=1.5)
fig.savefig(f'{CDIR}/c_dist.png', dpi=DPI, bbox_inches='tight'); plt.close()

# ── AA composition comparison ─────────────────────────────────────────────────
AAS = list('ACDEFGHIKLMNPQRSTVWY')
def aa_freq(df):
    fin = df[df.gen >= df.gen.max()-10]; c={a:0 for a in AAS}; t=0
    for seq in fin.sequence.dropna():
        for ch in str(seq):
            if ch in c: c[ch]+=1; t+=1
    return {a:c[a]/t*100 for a in AAS} if t else {a:0 for a in AAS}
fa, fb = aa_freq(A), aa_freq(B)
x = np.arange(len(AAS)); bw = 0.38
fig, ax = plt.subplots(figsize=FS_W)
ax.bar(x-bw/2, [fa[a] for a in AAS], bw, color=CA_s, alpha=0.82,
       label=LABEL_A, zorder=3)
ax.bar(x+bw/2, [fb[a] for a in AAS], bw, color=CB_s, alpha=0.82,
       label=LABEL_B, zorder=3)
ax.set_xticks(x); ax.set_xticklabels(AAS, fontsize=13)
ax.set_ylabel('Frequency in final population (%)')
ax.set_title('Amino Acid Composition — Final 10 Generations')
ax.legend(loc='upper right')
fig.tight_layout(pad=1.5)
fig.savefig(f'{CDIR}/c_aa.png', dpi=DPI, bbox_inches='tight'); plt.close()

print("Comparison plots saved.")

# ══════════════════════════════════════════════════════════════════════════════
# PPT PRIMITIVES
# ══════════════════════════════════════════════════════════════════════════════
NAVY  = RGBColor(0x1A,0x37,0x5A)
CA    = RGBColor(0x15,0x65,0xC0)
CB    = RGBColor(0x2E,0x7D,0x32)
WHITE = RGBColor(0xFF,0xFF,0xFF)
DGREY = RGBColor(0x22,0x22,0x22)
MGREY = RGBColor(0x60,0x60,0x60)
LGREY = RGBColor(0xF4,0xF4,0xF4)
RULE  = RGBColor(0xD0,0xD0,0xD0)
W, H  = 13.33, 7.5

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

def Rx(sl, x, y, w, h, fill=None):
    s = sl.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb=fill
    else: s.fill.background()
    s.line.fill.background()

def Tx(sl, text, x, y, w, h, size=14, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tb.word_wrap=True; tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color if color else DGREY

def I_fit(sl, path, x, y, max_w, max_h):
    if not os.path.exists(path): print(f"  ✗ missing: {path}"); return
    img=PILImage.open(path); iw,ih=img.size; r=ih/iw
    w=max_w; h=w*r
    if h>max_h: h=max_h; w=h/r
    ox=x+(max_w-w)/2; oy=y+(max_h-h)/2
    sl.shapes.add_picture(path,Inches(ox),Inches(oy),Inches(w),Inches(h))

def footer(sl):
    Rx(sl, 0, H-0.34, W, 0.34, fill=LGREY)
    Rx(sl, 0, H-0.34, W, 0.025, fill=RULE)
    Tx(sl, "National & Kapodistrian University of Athens   ·   Hellenic Pasteur Institute   ·   Evolutionary Genomics Group",
       0.3, H-0.3, W-0.6, 0.24, size=8.5, color=MGREY, align=PP_ALIGN.CENTER)

def chip(sl, x, y, text, color):
    cw = len(text)*0.092 + 0.4
    Rx(sl, x, y, cw, 0.30, fill=color)
    Tx(sl, text, x+0.1, y+0.04, cw-0.2, 0.22, size=11, bold=True, color=WHITE)
    return x+cw+0.16

# ── section divider slide ─────────────────────────────────────────────────────
def section_slide(title, subtitle, color):
    sl = prs.slides.add_slide(blank)
    Rx(sl, 0, 0, W, H, fill=color)
    Rx(sl, 0, H-0.36, W, 0.36, fill=RGBColor(0,0,0))  # slightly darker footer stripe
    Tx(sl, title,    0.6, 2.6, W-1.2, 1.0, size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    Tx(sl, subtitle, 0.6, 3.8, W-1.2, 0.6, size=18, color=RGBColor(0xCC,0xDD,0xFF) if color==NAVY else WHITE,
       align=PP_ALIGN.CENTER)

# ── graph slide ───────────────────────────────────────────────────────────────
IMG_TOP  = 1.08   # image starts here
IMG_AVAIL = H - IMG_TOP - 0.40   # space for image (leaving footer gap)

def graph_slide(title, img_path, run_color=None, run_label=None):
    sl = prs.slides.add_slide(blank)
    # header bar
    Rx(sl, 0, 0, W, 0.9, fill=NAVY)
    Tx(sl, title, 0.35, 0.1, W-0.7, 0.72, size=22, bold=True, color=WHITE)
    if run_color and run_label:
        chip(sl, 0.35, 0.08, run_label, run_color)
        Tx(sl, title, 0.35+len(run_label)*0.092+0.55, 0.1, W-0.7, 0.72, size=22, bold=True, color=WHITE)
    # full-width image
    I_fit(sl, img_path, 0.15, IMG_TOP, W-0.3, IMG_AVAIL)
    footer(sl)
    return sl

def graph_slide_clean(title, img_path, run_label=None, run_color=None):
    """Clean white header version for individual run plots."""
    sl = prs.slides.add_slide(blank)
    Tx(sl, title, 0.35, 0.1, W-0.7, 0.52, size=22, bold=True, color=NAVY)
    y_rule = 0.64
    if run_label and run_color:
        nx = chip(sl, 0.35, 0.66, run_label, run_color)
        y_rule = 1.02
    Rx(sl, 0.35, y_rule, W-0.7, 0.025, fill=RULE)
    img_start = y_rule + 0.06
    I_fit(sl, img_path, 0.15, img_start, W-0.3, H-img_start-0.42)
    footer(sl)
    return sl

# ══════════════════════════════════════════════════════════════════════════════
# BUILD SLIDES
# ══════════════════════════════════════════════════════════════════════════════

# 1. Title
sl = prs.slides.add_slide(blank)
Rx(sl, 0, 0, W, H, fill=NAVY)
Rx(sl, 0, H*0.55, W, H*0.45, fill=RGBColor(0x12,0x28,0x40))
Tx(sl, "Investigating novel AMPs\nusing machine learning,\nstructure prediction\nand in silico evolution",
   0.8, 0.7, W-1.6, 3.5, size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
Tx(sl, "MSc Bioinformatics – Computational Biology\nNational & Kapodistrian University of Athens\nApostolos Fysekidis",
   0.8, 4.5, W-1.6, 1.5, size=16, color=RGBColor(0xB0,0xC8,0xE8), align=PP_ALIGN.CENTER)
Rx(sl, 4.2, 6.3, 1.8, 0.06, fill=CA); Rx(sl, 7.3, 6.3, 1.8, 0.06, fill=CB)
footer(sl)

# ── SECTION 1: MACREL + ESM3 ──────────────────────────────────────────────────
section_slide("MACREL + ESM3", "Individual run analysis — 300 generations · 30,000 predictions", CA)

graph_slide_clean("Fitness Evolution",
    f"{DIR_A}/Evolution.png", LABEL_A, CA)
graph_slide_clean("Score Components",
    f"{DIR_A}/Score_components.png", LABEL_A, CA)
graph_slide_clean("Fitness Landscape",
    f"{DIR_A}/Fitness_landscape.png", LABEL_A, CA)
graph_slide_clean("Amino Acid Composition",
    f"{DIR_A}/AA_composition.png", LABEL_A, CA)
graph_slide_clean("Score Distribution",
    f"{DIR_A}/Score_distribution.png", LABEL_A, CA)
graph_slide_clean("Secondary Structure",
    f"{DIR_A}/Secondary_structures.png", LABEL_A, CA)
graph_slide_clean("Summary",
    f"{DIR_A}/Summary.png", LABEL_A, CA)

# ── SECTION 2: MACREL + ESM3 + PFES ──────────────────────────────────────────
section_slide("MACREL + ESM3 + PFES", "Individual run analysis — 300 generations · 30,000 predictions", CB)

graph_slide_clean("Fitness Evolution",
    f"{DIR_B}/Evolution.png", LABEL_B, CB)
graph_slide_clean("Score Components",
    f"{DIR_B}/Score_components.png", LABEL_B, CB)
graph_slide_clean("Fitness Landscape",
    f"{DIR_B}/Fitness_landscape.png", LABEL_B, CB)
graph_slide_clean("Amino Acid Composition",
    f"{DIR_B}/AA_composition.png", LABEL_B, CB)
graph_slide_clean("Score Distribution",
    f"{DIR_B}/Score_distribution.png", LABEL_B, CB)
graph_slide_clean("Secondary Structure",
    f"{DIR_B}/Secondary_structures.png", LABEL_B, CB)
graph_slide_clean("Summary",
    f"{DIR_B}/Summary.png", LABEL_B, CB)

# ── SECTION 3: Comparison ─────────────────────────────────────────────────────
section_slide("Comparison", f"{LABEL_A}  vs  {LABEL_B}", NAVY)

comp_plots = [
    ("Fitness Over Generations",                 f"{CDIR}/c_fitness.png"),
    ("Sequence Length Over Generations",         f"{CDIR}/c_length.png"),
    ("AMP Probability Over Generations",         f"{CDIR}/c_amp.png"),
    ("Hemolytic Proxy Over Generations",         f"{CDIR}/c_hemo.png"),
    ("pLDDT Over Generations",                   f"{CDIR}/c_plddt.png"),
    ("Amino Acid Composition — Comparison",      f"{CDIR}/c_aa.png"),
    ("Score Distribution — Final Generation",    f"{CDIR}/c_dist.png"),
]
for title, path in comp_plots:
    sl = prs.slides.add_slide(blank)
    Tx(sl, title, 0.35, 0.1, W-0.7, 0.52, size=22, bold=True, color=NAVY)
    nx = chip(sl, 0.35, 0.66, LABEL_A, CA)
    chip(sl, nx, 0.66, LABEL_B, CB)
    Rx(sl, 0.35, 1.02, W-0.7, 0.025, fill=RULE)
    I_fit(sl, path, 0.15, 1.08, W-0.3, H-1.08-0.42)
    footer(sl)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
n = len(prs.slides)
kb = os.path.getsize(OUT)//1024
print(f"Saved → {OUT}  ({kb} KB,  {n} slides)")
